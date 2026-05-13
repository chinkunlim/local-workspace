r"""
Phase 0c: MarkItDown 前置轉換
=================================
將非 PDF 格式（.pptx / .docx / .xlsx）轉換為統一的 Markdown 格式，
輸出 {file_id}_raw_extracted.md，介面與 p01a_engine.py 完全相同。

設計決策：
- 使用 Microsoft MarkItDown 作為轉換引擎（MIT License）
- PPTX：保留投影片結構（<!-- Slide number: N -->）+ 備注（### Notes:）
- DOCX：保留標題/段落層級
- XLSX：保留 Markdown 表格格式
- 圖片：從 PPTX 提取圖片 blob，存至與 raw_extracted.md 相同目錄，
         命名規則與 MarkItDown 一致（shape.name → re.sub(\W) + .jpg），
         讓 markdown 的 ![alt](filename.jpg) 引用直接生效
- 輸出 raw_extracted.md 標記為 IMMUTABLE，與 p1a_engine 介面一致
- 完成後將所有 PDF 專屬 phases (p0b, p0a, p1a, p1b_s, p1b, p1c, p1d) 設為 ⏭️

依賴：
  markitdown[pptx,docx,xlsx]（已加入 pyproject.toml）
  python-pptx（隨 markitdown[pptx] 自動安裝）
"""

import os
import re
import sys

# Internal Core Bootstrap
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "..")))
from core.utils.bootstrap import ensure_core_path as _bootstrap

_bootstrap(__file__)

from core import AtomicWriter, PipelineBase

# Phases that are PDF-specific and should be skipped for Office format files
_PDF_ONLY_PHASES = ["p0b", "p0a", "p1a", "p1b_s", "p1b", "p1c", "p1d"]

# Supported Office formats handled by this phase
_OFFICE_EXTENSIONS = {".pptx", ".docx", ".xlsx"}


class Phase0cMarkItDown(PipelineBase):
    """Phase 0c: MarkItDown conversion for Office format files."""

    def __init__(self) -> None:
        super().__init__(
            phase_key="p0c",
            phase_name="MarkItDown 前置轉換",
            skill_name="doc_parser",
        )

    def run(
        self,
        force: bool = False,
        subject: str = None,
        file_filter: str = None,
        single_mode: bool = False,
        resume_from: dict = None,
    ) -> None:
        self.log("🔄 啟動 Phase 0c：MarkItDown 前置轉換")

        # Inject Office-format files into state and mask PDF-only phases
        with self.state_manager._lock:
            for subj, files in self.state_manager.state.items():
                if subj == "_simple_":
                    continue
                for fname, record in files.items():
                    ext = os.path.splitext(fname)[1].lower()
                    if ext not in _OFFICE_EXTENSIONS:
                        continue
                    # Mark all PDF-specific phases as skipped for Office files
                    for p in _PDF_ONLY_PHASES:
                        if record.get(p) != "⏭️":
                            record[p] = "⏭️"
            self.state_manager._save_state()

        self.process_tasks(
            self._process_file,
            force=force,
            subject_filter=subject,
            file_filter=file_filter,
            single_mode=single_mode,
            resume_from=resume_from,
        )

    def _process_file(self, idx: int, task: dict, total: int) -> None:
        subj = task["subject"]
        fname = task["filename"]
        ext = os.path.splitext(fname)[1].lower()

        # Only handle Office formats in this phase
        if ext not in _OFFICE_EXTENSIONS:
            return

        file_id = os.path.splitext(fname)[0]
        inbox_path = os.path.join(self.dirs.get("inbox", ""), subj, fname)

        if not os.path.exists(inbox_path):
            self.log(f"⚠️ [{idx}/{total}] 找不到輸入檔：{inbox_path}", "warn")
            return

        self.log(f"📄 [{idx}/{total}] MarkItDown 轉換中：{fname}")

        # ── Run MarkItDown conversion ──
        try:
            from markitdown import MarkItDown

            md = MarkItDown(enable_plugins=False)
            result = md.convert(inbox_path)
            markdown_text = result.text_content or ""
        except ImportError:
            self.log("❌ markitdown 未安裝。請執行：uv add 'markitdown[pptx,docx,xlsx]'", "error")
            return
        except Exception as e:
            self.log(f"❌ MarkItDown 轉換失敗：{e}", "error")
            self.state_manager.update_task(subj, fname, self.phase_key, "❌", note_tag=str(e))
            return

        if not markdown_text.strip():
            self.log(f"⚠️ [{idx}/{total}] 轉換結果為空，跳過：{fname}", "warn")
            self.state_manager.update_task(subj, fname, self.phase_key, "❌", note_tag="空內容")
            return

        # ── Write raw_extracted.md (IMMUTABLE) ──
        output_dir = os.path.join(self.dirs["processed"], subj, file_id)
        os.makedirs(output_dir, exist_ok=True)
        raw_output_path = os.path.join(output_dir, f"{file_id}_raw_extracted.md")

        header = (
            f"<!-- {file_id}_raw_extracted.md — IMMUTABLE — DO NOT MODIFY —\n"
            f"     file_id: {file_id}\n"
            f"     source:  {fname}\n"
            f"     engine:  MarkItDown {ext.lstrip('.')}\n"
            "-->\n\n"
        )
        AtomicWriter.write_text(raw_output_path, header + markdown_text)

        char_count = len(markdown_text)
        self.log(f"✅ [{idx}/{total}] 轉換完成：{file_id} ({char_count:,} 字元)")

        # ── Extract embedded images (PPTX only) ──
        image_names: list[str] = []
        if ext == ".pptx":
            image_names = self._extract_pptx_images(inbox_path, output_dir)
            if image_names:
                self.log(f"🖼️  [{idx}/{total}] 提取圖片 {len(image_names)} 張")

        # ── Write figure_list.md ──
        figure_list_path = os.path.join(output_dir, "figure_list.md")
        if not os.path.exists(figure_list_path):
            rows = ""
            for img in image_names:
                rows += f"| {img} | — | — | — | — | PPTX embedded |\n"
            AtomicWriter.write_text(
                figure_list_path,
                (
                    "# Figure List\n\n"
                    "> 由 MarkItDown 轉換，圖片已提取至同目錄。\n\n"
                    "| 檔案名稱 | 頁碼 | 原始 Caption | VLM 描述 | 數據趨勢標籤 | 來源 |\n"
                    "| :--- | :---: | :--- | :--- | :---: | :---: |\n" + rows
                ),
            )

        # ── Update DAG state ──
        out_hash = self.state_manager.get_file_hash(raw_output_path)
        self.state_manager.update_task(
            subj,
            fname,
            self.phase_key,
            "✅",
            char_count=char_count,
            output_hash=out_hash,
        )

    def _extract_pptx_images(self, pptx_path: str, output_dir: str) -> list[str]:
        """Extract embedded images from a PPTX file.

        Uses python-pptx (installed with markitdown[pptx]) to read image
        blobs from every shape.  Files are saved to ``output_dir`` using
        the same filename convention as MarkItDown's _pptx_converter:

            re.sub(r"\\W", "", shape.name) + ".jpg"

        This ensures the ``![]()`` references already in raw_extracted.md
        resolve correctly without any post-processing.

        Returns a list of saved filenames (deduped).
        """
        try:
            import pptx as _pptx
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            self.log("⚠️ python-pptx 未找到，跳過圖片提取。", "warn")
            return []

        saved: list[str] = []
        seen: set[str] = set()

        def _save_shape_image(shape) -> None:  # type: ignore[no-untyped-def]
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                return
            try:
                image = shape.image
                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                dest = os.path.join(output_dir, filename)
                if filename not in seen:
                    with open(dest, "wb") as f:
                        f.write(image.blob)
                    seen.add(filename)
                    saved.append(filename)
            except Exception as exc:
                self.log(f"⚠️ 圖片提取失敗 ({shape.name}): {exc}", "warn")

        try:
            prs = _pptx.Presentation(pptx_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    _save_shape_image(shape)
                    # Recurse into group shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for subshape in shape.shapes:
                            _save_shape_image(subshape)
        except Exception as exc:
            self.log(f"❌ PPTX 圖片提取中斷：{exc}", "error")

        return saved


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Phase 0c: MarkItDown Office Conversion")
    parser.add_argument("--force", "-f", action="store_true")
    parser.add_argument("--subject", "-s", type=str)
    parser.add_argument("--file", type=str)
    args = parser.parse_args()
    Phase0cMarkItDown().run(force=args.force, subject=args.subject, file_filter=args.file)
